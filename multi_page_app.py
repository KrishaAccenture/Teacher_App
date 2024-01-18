import streamlit as st
import openai
import docx
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# Ensure that you have set up your st.secrets with the key "api_key"
openai.api_key = st.secrets["api_key"]

# Initialize session state variables for history
if 'history' not in st.session_state:
    st.session_state['history'] = []

# Function to update history
def update_history(session_details):
    st.session_state['history'].insert(0, session_details)

# Function to show history entry details
def show_history_entry_details(entry):
    st.write(f"Subject: {entry['user_input']['subject']}, Lesson Topic: {entry['user_input']['lesson_topic']}")
    st.download_button(label="Download Lesson Plan",
                       data=entry['lesson_plan_file_stream'],
                       file_name="lesson_plan.docx",
                       mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    st.download_button(label="Download PowerPoint Presentation",
                       data=entry['ppt_file_stream'],
                       file_name="presentation.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    st.download_button(label="Download Activity Sheets",
                       data=entry['activity_sheet_file_stream'],
                       file_name="activity_sheets.docx",
                       mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# Function to create a Word document
def create_word_document(content):
    doc = docx.Document()
    for paragraph in content.split('\n'):
        run = doc.add_paragraph().add_run(paragraph)
        font = run.font
        font.size = Pt(12)
    doc_stream = io.BytesIO()
    doc.save(doc_stream)
    doc_stream.seek(0)
    return doc_stream

# Function to create a PowerPoint presentation
# Function to create a PowerPoint presentation
def create_powerpoint(slides_content):
    prs = Presentation()
    for slide_content in slides_content:
        slide_layout = prs.slide_layouts[5]  # Using a title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set the title for the slide
        title_shape = slide.shapes.title
        title_shape.text = slide_content['title']
        
        # Create a textbox for the content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8.25)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = slide_content['content']
        p.font.size = PptPt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
    prs_stream = io.BytesIO()
    prs.save(prs_stream)
    prs_stream.seek(0)
    return prs_stream


# Main app functionality
def main():
    st.title('Lesson and Presentation Generator')

    subject = st.text_input("Subject", value="")
    year_group = st.text_input("Year Group", value="")
    lesson_topic = st.text_input("Lesson Topic", value="")
    number_of_lessons_required = st.number_input("Number of Lessons Required", min_value=1, value=1)
    ability_of_students = st.text_input("Ability of Students", value="")
    special_education_requirements = st.text_area("Special Education Requirements from Children", value="")
    additional_comments = st.text_area("Additional Comments", value="")

    # File uploader (optional for additional inputs)
    uploaded_files = st.file_uploader("Upload Supporting Files", accept_multiple_files=True, type=['pdf', 'docx', 'xlsx', 'csv', 'ppt', 'pptx'])

    # Submit button to generate materials
    if st.button('Click here to generate lesson plan, ppt, and activity sheets'):
        with st.spinner('Creating the lesson plan...'):
            # Placeholder for actual lesson plan content generation
            lesson_plan = "Generated lesson plan content"
            lesson_plan_file_stream = create_word_document(lesson_plan)
            
            # Placeholder for actual PowerPoint slides content generation
            ppt_content = "Generated PowerPoint content"
            ppt_file_stream = create_powerpoint([{"title": "Slide 1", "content": ppt_content}])
            
            # Placeholder for actual activity sheets content generation
            activity_sheet_content = "Generated activity sheet content"
            activity_sheet_file_stream = create_word_document(activity_sheet_content)
            
            # Update history with the new entry
            update_history({
                'user_input': {
                    'subject': subject,
                    'year_group': year_group,
                    'lesson_topic': lesson_topic,
                    'number_of_lessons_required': number_of_lessons_required,
                    'ability_of_students': ability_of_students,
                    'special_education_requirements': special_education_requirements,
                    'additional_comments': additional_comments
                },
                'lesson_plan': lesson_plan,
                'lesson_plan_file_stream': lesson_plan_file_stream,
                'ppt_content': ppt_content,
                'ppt_file_stream': ppt_file_stream,
                'activity_sheet_content': activity_sheet_content,
                'activity_sheet_file_stream': activity_sheet_file_stream
            })
            
            st.success('Materials ready for download!')

        # Display history in the sidebar
    with st.sidebar:
        st.header("History")
        for i, entry in enumerate(st.session_state['history']):
            if st.button(f"{entry['user_input']['subject']} - {entry['user_input']['lesson_topic']}", key=f"history_btn_{i}"):
                show_history_entry_details(i)

        if st.button('Start New', key='start_new'):
            # Reset the session state for user input
            for key in st.session_state.keys():
                del st.session_state[key]
            st.experimental_rerun()

if __name__ == "__main__":
    main()

